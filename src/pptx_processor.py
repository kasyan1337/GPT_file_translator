# src/pptx_processor.py

from pptx import Presentation
from src.document_processor import DocumentProcessor
from tqdm import tqdm
from src.utils import chunk_shapes


class PptxProcessor(DocumentProcessor):
    def process(self):
        presentation = Presentation(self.file_path)
        total_usage = {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0}

        # Collect all shapes with text
        shapes_with_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shapes_with_text.append((shape, shape.text))

        # Group shapes into chunks
        shape_chunks = chunk_shapes(
            shapes_with_text, max_tokens=7500, model="gpt-4-turbo"
        )

        for chunk in tqdm(shape_chunks, desc="Translating PPTX"):
            # Combine the text of the shapes in the chunk
            original_texts = [text for shape, text in chunk]
            original_text = "\n".join(original_texts)

            # Translate the combined text
            translated_text = self.openai_api.translate_text(self.prompt, original_text)
            if translated_text:
                # Split the translated text back into individual texts
                translated_texts = translated_text.split("\n")
                # Assign translated text back to the corresponding shapes
                for (shape, _), t in zip(chunk, translated_texts):
                    shape.text = t
                # Update token usage
                usage = self.openai_api.last_usage
                total_usage["prompt_tokens"] += usage["prompt_tokens"]
                total_usage["completion_tokens"] += usage["completion_tokens"]
                total_usage["total_tokens"] += usage["total_tokens"]
            else:
                print("Translation failed for a chunk.")

        presentation.save(self.output_path)

        # Print token usage for the file
        print(f"Token usage for {self.file_path}: {total_usage}")
        estimated_cost = self.openai_api.calculate_cost(total_usage)
        print(f"Estimated cost for {self.file_path}: ${estimated_cost:.4f}")
